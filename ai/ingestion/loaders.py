from __future__ import annotations

import re
from abc import ABC, abstractmethod
from pathlib import Path

import fitz
from docx import Document as WordDocument
from pptx import Presentation

from ai.exceptions import (
    DocumentParseError,
    EmptyDocumentError,
    UnsupportedDocumentTypeError,
)
from ai.models import ParsedDocument, ParsedSection

from typing import Any

from ai.ingestion.ocr import (
    FallbackOCRService,
    OCRResult,
    OCRServiceProtocol,
    OpenAICompatibleVisionOCRService,
    PaddleOCRService,
)


PARSER_VERSION = "document-parser-v1"

def normalize_text(text:str) -> str:
    """统一换行和行尾空白，同时保留段落边界。"""

    text = text.replace("\r\n", "\n").replace("\r", "\n")
    lines = [line.rstrip() for line in text.splitlines()]

    normalized : list[str] = []
    previous_blank = False

    for line in lines:
        blank = not line.strip()

        if blank:
            if normalized and not previous_blank:
                normalized.append("")
        else:
            normalized.append(line.strip())

        previous_blank = blank

    return "\n".join(normalized).strip()

def pixmap_to_ndarray(pixmap: fitz.Pixmap) -> Any:
    """将 PyMuPDF Pixmap 转为 PaddleOCR 接受的 RGB ndarray。"""

    try:
        import numpy as np
    except ImportError as exc:
        raise DocumentParseError("OCR 需要 numpy。") from exc

    channels = pixmap.n
    array = np.frombuffer(
        pixmap.samples,
        dtype=np.uint8,
    ).reshape(pixmap.height, pixmap.width, channels)

    if channels == 4:
        array = array[:, :, :3]

    return array.copy()

class BaseDocumentParser(ABC):
    """格式解析器的统一接口。"""

    file_type: str
    supported_suffixes: frozenset[str]

    @abstractmethod
    def parse(self, path: Path) -> ParsedDocument:
        raise NotImplementedError

    def _build_document(
        self,
        path: Path,
        sections: list[ParsedSection],
    ) -> ParsedDocument:
        if not sections:
            raise EmptyDocumentError(f"没有从文件中提取到文本：{path.name}")

        return ParsedDocument(
            source_path=path,
            file_type=self.file_type,
            parser_version=PARSER_VERSION,
            sections=tuple(sections),
        )


class PDFParser(BaseDocumentParser):
    file_type = "pdf"
    supported_suffixes = frozenset({".pdf"})

    def __init__(
        self,
        *,
        ocr_service: OCRServiceProtocol | None = None,
        ocr_enabled: bool = False,
        ocr_dpi: int = 300,
        min_native_characters: int = 30,
    ) -> None:
        self.ocr_service = ocr_service
        self.ocr_enabled = ocr_enabled
        self.ocr_dpi = ocr_dpi
        self.min_native_characters = min_native_characters

    def parse(self, path: Path) -> ParsedDocument:
        sections: list[ParsedSection] = []

        try:
            with fitz.open(path) as pdf:
                for page_index, page in enumerate(pdf):
                    page_number = page_index + 1
                    native_text = normalize_text(page.get_text("text"))

                    content = native_text
                    extraction_method = "native"
                    ocr_confidence: float | None = None
                    ocr_line_count = 0

                    needs_ocr = (
                        self.ocr_enabled
                        and len(native_text) < self.min_native_characters
                    )

                    if needs_ocr:
                        try:
                            ocr_result = self._ocr_page(page)
                        except DocumentParseError:
                            content = "[扫描版 PDF：此页等待 OCR 识别]"
                            extraction_method = "ocr_pending"
                            ocr_result = None

                        # OCR 必须比原生提取结果更有信息才替换。
                        if ocr_result is not None and len(ocr_result.text.strip()) > len(native_text):
                            content = normalize_text(ocr_result.text)
                            extraction_method = "ocr"
                            ocr_confidence = ocr_result.average_confidence
                            ocr_line_count = ocr_result.line_count

                    if not content:
                        continue

                    sections.append(
                        ParsedSection(
                            content=content,
                            source_path=path,
                            source_name=path.name,
                            file_type=self.file_type,
                            page_start=page_number,
                            page_end=page_number,
                            location_label=f"第 {page_number} 页",
                            metadata={
                                "pdf_page_index": page_index,
                                "parser_version": PARSER_VERSION,
                                "extraction_method": extraction_method,
                                "ocr_confidence": ocr_confidence,
                                "ocr_line_count": ocr_line_count,
                            },
                        )
                    )
        except DocumentParseError:
            raise
        except Exception as exc:
            raise DocumentParseError(
                f"PDF 解析失败：{path.name}：{exc}"
            ) from exc

        return self._build_document(path, sections)

    def _ocr_page(self, page: fitz.Page) -> OCRResult:
        if self.ocr_service is None:
            raise DocumentParseError(
                "检测到扫描 PDF，但没有配置 OCR 服务。"
            )

        scale = self.ocr_dpi / 72.0
        matrix = fitz.Matrix(scale, scale)
        pixmap = page.get_pixmap(
            matrix=matrix,
            colorspace=fitz.csRGB,
            alpha=False,
        )
        image = pixmap_to_ndarray(pixmap)

        return self.ocr_service.recognize(image)

class ImageParser(BaseDocumentParser):
    file_type = "image"
    supported_suffixes = frozenset({
        ".png",
        ".jpg",
        ".jpeg",
        ".bmp",
        ".tif",
        ".tiff",
        ".webp",
    })

    def __init__(self, ocr_service: OCRServiceProtocol) -> None:
        self.ocr_service = ocr_service

    def parse(self, path: Path) -> ParsedDocument:
        try:
            result = self.ocr_service.recognize(path)
            content = normalize_text(result.text)

            if not content:
                raise EmptyDocumentError(
                    f"没有从图片中识别到文字：{path.name}"
                )

            section = ParsedSection(
                content=content,
                source_path=path,
                source_name=path.name,
                file_type=self.file_type,
                location_label="图片",
                metadata={
                    "image_number": 1,
                    "parser_version": PARSER_VERSION,
                    "extraction_method": "ocr",
                    "ocr_confidence": result.average_confidence,
                    "ocr_line_count": result.line_count,
                },
            )

            return self._build_document(path, [section])
        except (DocumentParseError, EmptyDocumentError):
            raise
        except Exception as exc:
            raise DocumentParseError(
                f"图片 OCR 失败：{path.name}：{exc}"
            ) from exc

class DOCXParser(BaseDocumentParser):
    file_type = "docx"
    supported_suffixes = frozenset({".docx"})

    def parse(self, path: Path) -> ParsedDocument:
        sections: list[ParsedSection] = []

        try:
            document = WordDocument(path)

            for paragraph_index, paragraph in enumerate(document.paragraphs):
                content = normalize_text(paragraph.text)

                if not content:
                    continue

                paragraph_number = paragraph_index + 1
                style_name = paragraph.style.name if paragraph.style else ""
                is_heading = style_name.lower().startswith("heading")

                sections.append(
                    ParsedSection(
                        content=content,
                        source_path=path,
                        source_name=path.name,
                        file_type=self.file_type,
                        location_label=f"第 {paragraph_number} 段",
                        section_title=content if is_heading else "",
                        metadata={
                            "paragraph_number": paragraph_number,
                            "style_name": style_name,
                            "parser_version": PARSER_VERSION,
                        },
                    )
                )
        except Exception as exc:
            raise DocumentParseError(
                f"DOCX 解析失败：{path.name}：{exc}"
            ) from exc

        return self._build_document(path, sections)


class PPTXParser(BaseDocumentParser):
    file_type = "pptx"
    supported_suffixes = frozenset({".pptx"})

    def parse(self, path: Path) -> ParsedDocument:
        sections: list[ParsedSection] = []

        try:
            presentation = Presentation(path)

            for slide_index, slide in enumerate(presentation.slides):
                slide_number = slide_index + 1
                text_parts: list[str] = []

                for shape in slide.shapes:
                    if not hasattr(shape, "text"):
                        continue

                    text = normalize_text(shape.text)
                    if text:
                        text_parts.append(text)

                content = normalize_text("\n".join(text_parts))

                if not content:
                    continue

                title = ""
                if slide.shapes.title is not None:
                    title = normalize_text(slide.shapes.title.text)

                sections.append(
                    ParsedSection(
                        content=content,
                        source_path=path,
                        source_name=path.name,
                        file_type=self.file_type,
                        location_label=f"第 {slide_number} 张幻灯片",
                        section_title=title,
                        metadata={
                            "slide_number": slide_number,
                            "parser_version": PARSER_VERSION,
                        },
                    )
                )
        except Exception as exc:
            raise DocumentParseError(
                f"PPTX 解析失败：{path.name}：{exc}"
            ) from exc

        return self._build_document(path, sections)

class MarkdownParser(BaseDocumentParser):
    file_type = "markdown"
    supported_suffixes = frozenset({".md", ".markdown"})

    _heading_pattern = re.compile(r"^(#{1,6})\s+(.+?)\s*$")

    def parse(self, path: Path) -> ParsedDocument:
        try:
            text = path.read_text(encoding="utf-8-sig")
        except UnicodeError as exc:
            raise DocumentParseError(
                f"Markdown 不是有效的 UTF-8 文件：{path.name}"
            ) from exc
        except OSError as exc:
            raise DocumentParseError(
                f"无法读取 Markdown：{path.name}：{exc}"
            ) from exc

        lines = text.replace("\r\n", "\n").replace("\r", "\n").splitlines()
        sections: list[ParsedSection] = []

        current_lines: list[str] = []
        current_title = ""
        current_heading_level: int | None = None
        section_start = 1

        def append_current(end_line: int) -> None:
            content = normalize_text("\n".join(current_lines))

            if not content:
                return

            sections.append(
                ParsedSection(
                    content=content,
                    source_path=path,
                    source_name=path.name,
                    file_type=self.file_type,
                    line_start=section_start,
                    line_end=end_line,
                    location_label=f"第 {section_start}–{end_line} 行",
                    section_title=current_title,
                    metadata={
                        "heading_level": current_heading_level,
                        "parser_version": PARSER_VERSION,
                    },
                )
            )

        for line_number, line in enumerate(lines, start=1):
            heading = self._heading_pattern.match(line)

            if heading and current_lines:
                append_current(line_number - 1)
                current_lines = []
                section_start = line_number

            if heading:
                current_heading_level = len(heading.group(1))
                current_title = heading.group(2).strip()

            current_lines.append(line)

        if current_lines:
            append_current(len(lines))

        return self._build_document(path, sections)



class TXTParser(BaseDocumentParser):
    file_type = "txt"
    supported_suffixes = frozenset({".txt"})

    def parse(self, path: Path) -> ParsedDocument:
        try:
            text = path.read_text(encoding="utf-8-sig")
        except UnicodeError as exc:
            raise DocumentParseError(
                f"TXT 不是有效的 UTF-8 文件：{path.name}"
            ) from exc
        except OSError as exc:
            raise DocumentParseError(
                f"无法读取 TXT：{path.name}：{exc}"
            ) from exc

        lines = text.replace("\r\n", "\n").replace("\r", "\n").splitlines()
        sections: list[ParsedSection] = []

        paragraph_lines: list[str] = []
        paragraph_start: int | None = None

        def append_paragraph(end_line: int) -> None:
            nonlocal paragraph_lines, paragraph_start

            if paragraph_start is None:
                return

            content = normalize_text("\n".join(paragraph_lines))

            if content:
                sections.append(
                    ParsedSection(
                        content=content,
                        source_path=path,
                        source_name=path.name,
                        file_type=self.file_type,
                        line_start=paragraph_start,
                        line_end=end_line,
                        location_label=f"第 {paragraph_start}–{end_line} 行",
                        metadata={
                            "parser_version": PARSER_VERSION,
                        },
                    )
                )

            paragraph_lines = []
            paragraph_start = None

        for line_number, line in enumerate(lines, start=1):
            if line.strip():
                if paragraph_start is None:
                    paragraph_start = line_number
                paragraph_lines.append(line)
            else:
                append_paragraph(line_number - 1)

        append_paragraph(len(lines))

        return self._build_document(path, sections)



class DocumentParserRegistry:
    """根据扩展名选择解析器。"""

    def __init__(
        self,
        *,
        ocr_service: OCRServiceProtocol | None = None,
        ocr_enabled: bool = False,
        ocr_language: str = "ch",
        ocr_device: str = "cpu",
        ocr_dpi: int = 300,
        ocr_min_native_characters: int = 30,
        ocr_min_confidence: float = 0.5,
        ocr_detection_model_name: str = "PP-OCRv5_mobile_det",
        ocr_detection_model_dir: Path | None = None,
        ocr_recognition_model_name: str = "PP-OCRv5_mobile_rec",
        ocr_recognition_model_dir: Path | None = None,
        ocr_use_doc_orientation: bool = False,
        ocr_use_textline_orientation: bool = False,
        ocr_enable_mkldnn: bool = False,
    ) -> None:
        service = ocr_service

        if service is None and ocr_enabled:
            service = PaddleOCRService(
                language=ocr_language,
                device=ocr_device,
                min_confidence=ocr_min_confidence,
                detection_model_name=ocr_detection_model_name,
                detection_model_dir=ocr_detection_model_dir,
                recognition_model_name=ocr_recognition_model_name,
                recognition_model_dir=ocr_recognition_model_dir,
                use_doc_orientation=ocr_use_doc_orientation,
                use_textline_orientation=ocr_use_textline_orientation,
                enable_mkldnn=ocr_enable_mkldnn,
            )

        parsers: list[BaseDocumentParser] = [
            PDFParser(
                ocr_service=service,
                ocr_enabled=ocr_enabled,
                ocr_dpi=ocr_dpi,
                min_native_characters=ocr_min_native_characters,
            ),
            DOCXParser(),
            PPTXParser(),
            MarkdownParser(),
            TXTParser(),
        ]

        if service is not None:
            parsers.append(ImageParser(service))

        self._parsers = {
            suffix: parser
            for parser in parsers
            for suffix in parser.supported_suffixes
        }

    @property
    def supported_suffixes(self) -> frozenset[str]:
        return frozenset(self._parsers)

    def parse(self, path: Path | str) -> ParsedDocument:
        source = Path(path).expanduser().resolve()

        if not source.exists():
            raise DocumentParseError(f"文件不存在：{source}")

        if not source.is_file():
            raise DocumentParseError(f"目标不是文件：{source}")

        suffix = source.suffix.lower()
        parser = self._parsers.get(suffix)

        if parser is None:
            raise UnsupportedDocumentTypeError(
                f"暂不支持 {suffix or '无扩展名'} 文件；"
                f"支持：{', '.join(sorted(self.supported_suffixes))}"
            )

        return parser.parse(source)

def create_document_parser_registry() -> DocumentParserRegistry:
    """Create a parser registry from the current AI settings."""
    from ai.config import get_ai_settings

    settings = get_ai_settings()

    local_service: OCRServiceProtocol | None = None
    if settings.ocr_enabled:
        # With model directories unset, PaddleOCR downloads and caches its
        # official CPU weights on first use.
        local_service = PaddleOCRService(
            language=settings.ocr_language,
            device=settings.ocr_device,
            min_confidence=settings.ocr_min_confidence,
            detection_model_name=settings.ocr_detection_model_name,
            detection_model_dir=settings.ocr_detection_model_dir,
            recognition_model_name=settings.ocr_recognition_model_name,
            recognition_model_dir=settings.ocr_recognition_model_dir,
            use_doc_orientation=settings.ocr_use_doc_orientation,
            use_textline_orientation=settings.ocr_use_textline_orientation,
            enable_mkldnn=settings.ocr_enable_mkldnn,
        )
    service = local_service
    if settings.ocr_api_enabled:
        remote_service = OpenAICompatibleVisionOCRService(
            api_key=settings.ocr_api_key,
            base_url=str(settings.ocr_api_base_url),
            model=settings.ocr_api_model,
            timeout_seconds=settings.ocr_api_timeout_seconds,
        )
        service = FallbackOCRService(remote_service, local_service) if local_service else remote_service

    return DocumentParserRegistry(
        ocr_service=service,
        ocr_enabled=service is not None,
        ocr_language=settings.ocr_language,
        ocr_device=settings.ocr_device,
        ocr_dpi=settings.ocr_dpi,
        ocr_min_native_characters=settings.ocr_min_native_characters,
        ocr_min_confidence=settings.ocr_min_confidence,
        ocr_detection_model_name=settings.ocr_detection_model_name,
        ocr_detection_model_dir=None,
        ocr_recognition_model_name=settings.ocr_recognition_model_name,
        ocr_recognition_model_dir=None,
        ocr_use_doc_orientation=settings.ocr_use_doc_orientation,
        ocr_use_textline_orientation=settings.ocr_use_textline_orientation,
        ocr_enable_mkldnn=settings.ocr_enable_mkldnn,
    )
