from __future__ import annotations

import fitz
import pytest
from docx import Document as WordDocument
from pptx import Presentation

from typing import Any

from ai.ingestion.ocr import OCRResult

from ai.exceptions import (
    EmptyDocumentError,
    UnsupportedDocumentTypeError,
)
from ai.ingestion import DocumentParserRegistry


@pytest.fixture
def registry() -> DocumentParserRegistry:
    return DocumentParserRegistry()

class FakeOCRService:
    def __init__(self, text: str = "扫描识别文本") -> None:
        self.text = text
        self.call_count = 0

    def recognize(self, image: str | object | Any) -> OCRResult:
        self.call_count += 1
        return OCRResult(
            text=self.text,
            average_confidence=0.98,
            line_count=1,
        )

def test_scanned_pdf_uses_ocr_and_preserves_page_number(tmp_path) -> None:
    path = tmp_path / "扫描教材.pdf"

    pdf = fitz.open()
    pdf.new_page()
    pdf.new_page()
    pdf.save(path)
    pdf.close()

    fake_ocr = FakeOCRService("极限描述函数的变化趋势。")
    registry = DocumentParserRegistry(
        ocr_service=fake_ocr,
        ocr_enabled=True,
    )

    parsed = registry.parse(path)

    assert len(parsed.sections) == 2
    assert fake_ocr.call_count == 2

    first = parsed.sections[0]
    second = parsed.sections[1]

    assert first.page_start == 1
    assert first.location_label == "第 1 页"
    assert first.metadata["extraction_method"] == "ocr"
    assert first.metadata["ocr_confidence"] == 0.98

    assert second.page_start == 2
    assert second.location_label == "第 2 页"

    documents = parsed.to_langchain_documents()
    assert documents[0].metadata["page_start"] == 1
    assert documents[1].metadata["page_start"] == 2

def test_text_pdf_does_not_use_ocr_when_text_is_sufficient(
    tmp_path,
) -> None:
    path = tmp_path / "文本教材.pdf"

    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text(
        (72, 72),
        "This is a native PDF page with enough searchable text content.",
    )
    pdf.save(path)
    pdf.close()

    fake_ocr = FakeOCRService("不应该使用")
    registry = DocumentParserRegistry(
        ocr_service=fake_ocr,
        ocr_enabled=True,
        ocr_min_native_characters=30,
    )

    parsed = registry.parse(path)

    assert fake_ocr.call_count == 0
    assert parsed.sections[0].metadata["extraction_method"] == "native"

def test_image_uses_ocr_without_fake_page_number(tmp_path) -> None:
    path = tmp_path / "课堂板书.png"
    path.write_bytes(b"fake image content")

    fake_ocr = FakeOCRService("导数表示函数的瞬时变化率。")
    registry = DocumentParserRegistry(
        ocr_service=fake_ocr,
        ocr_enabled=True,
    )

    parsed = registry.parse(path)

    assert parsed.file_type == "image"
    assert len(parsed.sections) == 1

    section = parsed.sections[0]
    assert section.content == "导数表示函数的瞬时变化率。"
    assert section.location_label == "图片"
    assert section.page_start is None
    assert section.metadata["extraction_method"] == "ocr"

def test_parse_pdf_preserves_one_based_page_numbers(
    tmp_path,
    registry: DocumentParserRegistry,
) -> None:
    path = tmp_path / "高等数学.pdf"

    pdf = fitz.open()
    first_page = pdf.new_page()
    first_page.insert_text((72, 72), "Limit definition")

    second_page = pdf.new_page()
    second_page.insert_text((72, 72), "Continuity definition")

    pdf.save(path)
    pdf.close()

    parsed = registry.parse(path)

    assert parsed.file_type == "pdf"
    assert len(parsed.sections) == 2

    assert parsed.sections[0].page_start == 1
    assert parsed.sections[0].page_end == 1
    assert parsed.sections[0].location_label == "第 1 页"

    assert parsed.sections[1].page_start == 2
    assert parsed.sections[1].location_label == "第 2 页"

    documents = parsed.to_langchain_documents()
    assert documents[1].metadata["page_start"] == 2
    assert documents[1].metadata["source_name"] == "高等数学.pdf"


def test_parse_docx_preserves_paragraph_positions(
    tmp_path,
    registry: DocumentParserRegistry,
) -> None:
    path = tmp_path / "讲义.docx"

    document = WordDocument()
    document.add_heading("函数极限", level=1)
    document.add_paragraph("函数极限描述函数值的变化趋势。")
    document.save(path)

    parsed = registry.parse(path)

    assert parsed.file_type == "docx"
    assert len(parsed.sections) == 2
    assert parsed.sections[0].section_title == "函数极限"
    assert parsed.sections[1].metadata["paragraph_number"] == 2


def test_parse_pptx_preserves_slide_numbers(
    tmp_path,
    registry: DocumentParserRegistry,
) -> None:
    path = tmp_path / "课程.pptx"

    presentation = Presentation()
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[1]
    )
    slide.shapes.title.text = "导数"
    slide.placeholders[1].text = "导数描述函数的瞬时变化率。"
    presentation.save(path)

    parsed = registry.parse(path)

    assert parsed.file_type == "pptx"
    assert len(parsed.sections) == 1
    assert parsed.sections[0].section_title == "导数"
    assert parsed.sections[0].metadata["slide_number"] == 1
    assert parsed.sections[0].location_label == "第 1 张幻灯片"
    assert parsed.sections[0].page_start is None


def test_parse_markdown_preserves_heading_and_lines(
    tmp_path,
    registry: DocumentParserRegistry,
) -> None:
    path = tmp_path / "笔记.md"
    path.write_text(
        "# 极限\n"
        "极限描述变化趋势。\n"
        "\n"
        "## 连续\n"
        "连续性与极限有关。\n",
        encoding="utf-8",
    )

    parsed = registry.parse(path)

    assert len(parsed.sections) == 2
    assert parsed.sections[0].section_title == "极限"
    assert parsed.sections[0].line_start == 1
    assert parsed.sections[0].line_end == 3

    assert parsed.sections[1].section_title == "连续"
    assert parsed.sections[1].line_start == 4
    assert parsed.sections[1].line_end == 5


def test_parse_txt_preserves_paragraph_line_ranges(
    tmp_path,
    registry: DocumentParserRegistry,
) -> None:
    path = tmp_path / "课堂记录.txt"
    path.write_text(
        "第一段第一行\n"
        "第一段第二行\n"
        "\n"
        "第二段\n",
        encoding="utf-8",
    )

    parsed = registry.parse(path)

    assert len(parsed.sections) == 2
    assert parsed.sections[0].line_start == 1
    assert parsed.sections[0].line_end == 2
    assert parsed.sections[1].line_start == 4
    assert parsed.sections[1].line_end == 4


def test_rejects_unsupported_file_type(
    tmp_path,
    registry: DocumentParserRegistry,
) -> None:
    path = tmp_path / "data.csv"
    path.write_text("a,b", encoding="utf-8")

    with pytest.raises(UnsupportedDocumentTypeError):
        registry.parse(path)



def test_rejects_empty_text_document(
    tmp_path,
    registry: DocumentParserRegistry,
) -> None:
    path = tmp_path / "empty.txt"
    path.write_text("\n\n", encoding="utf-8")

    with pytest.raises(EmptyDocumentError):
        registry.parse(path)